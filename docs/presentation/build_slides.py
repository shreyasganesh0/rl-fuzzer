#!/usr/bin/env python3
"""Build Experiment 3 weekly review presentation as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

import os, sys

# ── Colours ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG2   = RGBColor(0x16, 0x21, 0x3E)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)
RED        = RGBColor(0xEF, 0x47, 0x6F)
ORANGE     = RGBColor(0xFF, 0xD1, 0x66)
GREY       = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_GREY = RGBColor(0xDD, 0xDD, 0xDD)
CHART_DQN  = RGBColor(0x00, 0xB4, 0xD8)
CHART_BAN  = RGBColor(0xFF, 0xD1, 0x66)
CHART_M10  = RGBColor(0x90, 0xE0, 0xEF)
CHART_BL   = RGBColor(0xEF, 0x47, 0x6F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=Pt(6),
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p

def add_bullet(tf, text, font_size=16, color=WHITE, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = Pt(4)
    return p

def title_bar(slide, title, subtitle=None):
    """Dark accent bar at top with title."""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(1.3))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG2
    shape.line.fill.background()

    add_text(slide, Inches(0.7), Inches(0.2), Inches(11), Inches(0.6),
             title, font_size=32, color=ACCENT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.75), Inches(11), Inches(0.4),
                 subtitle, font_size=16, color=GREY)

def section_slide(slide, number, title, subtitle=""):
    set_slide_bg(slide, DARK_BG2)
    add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             f"0{number}" if number < 10 else str(number),
             font_size=72, color=ACCENT, bold=True, alignment=PP_ALIGN.LEFT)
    add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.0),
             title, font_size=40, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(1), Inches(4.2), Inches(10), Inches(0.8),
                 subtitle, font_size=20, color=GREY)


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT,
              cell_color=WHITE, font_size=13):
    """data = list of lists (first row is header)."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for col_idx in range(cols):
        for row_idx in range(rows):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = DARK_BG
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    paragraph.font.color.rgb = cell_color
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table


# ==========================================================================
# SLIDE 1 — Title
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Experiment 3: Differential-Informed",
         font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_text(slide, Inches(1), Inches(2.4), Inches(11), Inches(1.0),
         "RL-Guided Fuzzing (Model M3_0)",
         font_size=40, color=ACCENT, bold=True, alignment=PP_ALIGN.LEFT)

add_text(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.5),
         "Weekly Review  |  April 4, 2026",
         font_size=20, color=GREY)
add_text(slide, Inches(1), Inches(4.6), Inches(10), Inches(0.5),
         "Shreyas Ganesh",
         font_size=22, color=WHITE, bold=True)

# accent line
shape = slide.shapes.add_shape(1, Inches(1), Inches(3.6), Inches(4), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()


# ==========================================================================
# SLIDE 2 — Agenda
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Agenda")

items = [
    "1.  Motivation & Research Question",
    "2.  Experiment Design (Differential Fuzzing Pipeline)",
    "3.  Phase 1-2: Target Construction & Telemetry Collection",
    "4.  Phase 3: Differential Analysis & Feature Derivation",
    "5.  Phase 4: M3_0 Implementation",
    "6.  Phase 5: Training & Evaluation Protocol",
    "7.  Results",
    "8.  Analysis & Discussion",
    "9.  Next Steps",
]
tf = add_text(slide, Inches(1), Inches(1.7), Inches(10), Inches(5),
              items[0], font_size=22, color=WHITE)
for item in items[1:]:
    add_para(tf, item, font_size=22, color=WHITE, space_before=Pt(14))


# ==========================================================================
# SLIDE 3 — Motivation
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 1, "Motivation", "Why differential-informed features?")


# ==========================================================================
# SLIDE 4 — The Problem
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "The Problem: State Representation for RL Fuzzing")

tf = add_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5),
              "Prior RL models used hand-designed features:", font_size=18, color=WHITE, bold=True)
add_bullet(tf, "M0_0: 3 dims (coverage, new_edges, crashes) \u2014 too sparse", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "M1_0: 12 dims (edge stability distribution) \u2014 intuition-based", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "M1_1: 13 dims (visited-edge tracking) \u2014 high overhead", font_size=16, color=LIGHT_GREY)
add_para(tf, "", font_size=10)
add_para(tf, "Core weakness:", font_size=18, color=RED, bold=True)
add_bullet(tf, "No empirical evidence these features correlate with bug-finding", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "All designed by intuition about what \"should\" matter", font_size=16, color=LIGHT_GREY)
add_para(tf, "", font_size=10)
add_para(tf, "The differential insight:", font_size=18, color=GREEN, bold=True)
add_bullet(tf, "Fuzz buggy vs. fixed versions of the same software", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "Measure which execution metrics actually diverge", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "Use those as the RL state representation", font_size=16, color=LIGHT_GREY)

# Right side: model progression table
data = [
    ["Model", "Dims", "Feature Source", "Validated?"],
    ["M0_0",  "3",    "Manual",         "No"],
    ["M1_0",  "12",   "Manual",         "No"],
    ["M1_1",  "13",   "Manual",         "No"],
    ["M3_0",  "13",   "Differential",   "Yes"],
]
add_table(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.5),
          5, 4, data, font_size=14)


# ==========================================================================
# SLIDE 5 — Research Questions
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Research Questions")

questions = [
    ("RQ1", "Can RL features derived from differential analysis outperform\n"
            "hand-designed features for mutation selection?"),
    ("RQ2", "Do differential features transfer to unseen targets\n"
            "(different CVE, same codebase)?"),
    ("RQ3", "DQN vs. contextual bandit \u2014 which learns better policies\n"
            "from these features?"),
    ("RQ4", "What is the current gap between RL-guided and vanilla AFL++,\n"
            "and what causes it?"),
]
y = 1.7
for tag, text in questions:
    add_text(slide, Inches(0.7), Inches(y), Inches(1.2), Inches(0.5),
             tag, font_size=24, color=ACCENT, bold=True)
    add_text(slide, Inches(2.0), Inches(y), Inches(10), Inches(0.8),
             text, font_size=18, color=WHITE)
    y += 1.2


# ==========================================================================
# SLIDE 6 — Section: Experiment Design
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 2, "Experiment Design",
              "5-phase differential fuzzing pipeline")


# ==========================================================================
# SLIDE 7 — Architecture
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "System Architecture", "AFL++ + Custom Mutator + RL Server via Shared Memory IPC")

# Architecture diagram as text
arch_text = (
    "\u250C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510"
    "     mmap'd SHM (128 bytes)     "
    "\u250C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n"
    "\u2502   AFL++ Fuzzer    \u2502"
    " \u25C4\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u25BA "
    "\u2502  RL Server (Py)   \u2502\n"
    "\u2502  + Custom Mutator  \u2502"
    "  13 features \u2192 action (0\u201346)   "
    "\u2502  DQN / Bandit      \u2502\n"
    "\u2502    (.so plugin)    \u2502"
    "  lock-free atomic protocol    "
    "\u2502  [128,128,64] net  \u2502\n"
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518"
    "                                "
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518"
)
add_text(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(2.5),
         arch_text, font_size=14, color=ACCENT2, font_name="Consolas")

# Key details below
tf = add_text(slide, Inches(0.7), Inches(3.8), Inches(5.5), Inches(3.5),
              "Mutator (C, ~480 lines):", font_size=18, color=ACCENT, bold=True)
add_bullet(tf, "Computes 13 features per execution", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Cumulative bitmap (max-merge of trace_bits)", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Edge heat, entropy, timing, velocity", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Writes state, waits for action via atomic SHM", font_size=15, color=LIGHT_GREY)

tf2 = add_text(slide, Inches(6.8), Inches(3.8), Inches(5.5), Inches(3.5),
               "RL Server (Python):", font_size=18, color=ACCENT, bold=True)
add_bullet(tf2, "Reads 13-dim state, selects from 47 mutations", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "DQN: Double DQN, \u03b5-greedy, replay buffer 100K", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Bandit: Thompson sampling, 2-head network", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Reward: new_edges + 10 \u00d7 crashes", font_size=15, color=LIGHT_GREY)


# ==========================================================================
# SLIDE 8 — Pipeline Overview
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "5-Phase Pipeline")

phases = [
    ("Phase 1", "Build Targets", "4 libxml2 binaries\n(2 CVEs \u00d7 buggy/fixed)\nAFL++ instrumented + ASAN", ACCENT),
    ("Phase 2", "Telemetry", "24 campaigns (~10h)\n40K bitmap snapshots\n3.4 GB telemetry data", ACCENT2),
    ("Phase 3", "Differential\nAnalysis", "Mann-Whitney U test\nA12 effect size ranking\n13 features selected", GREEN),
    ("Phase 4", "Implement\nM3_0", "C mutator rewrite\nPython model rewrite\n128-byte SHM protocol", ORANGE),
    ("Phase 5", "Train & Eval", "500K steps training\n5 eval runs \u00d7 4 variants\n2 targets (in-dist + transfer)", RED),
]

x_start = 0.3
box_w = 2.35
gap = 0.2
for i, (phase, title, desc, color) in enumerate(phases):
    x = x_start + i * (box_w + gap)
    # Box
    shape = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(box_w), Inches(4.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    # Phase label
    add_text(slide, Inches(x + 0.15), Inches(1.9), Inches(box_w - 0.3), Inches(0.4),
             phase, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text(slide, Inches(x + 0.15), Inches(2.3), Inches(box_w - 0.3), Inches(0.8),
             title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text(slide, Inches(x + 0.15), Inches(3.3), Inches(box_w - 0.3), Inches(2.5),
             desc, font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(phases) - 1:
        add_text(slide, Inches(x + box_w - 0.05), Inches(3.3), Inches(gap + 0.2), Inches(0.5),
                 "\u25B6", font_size=20, color=GREY, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# SLIDE 9 — Section: Targets & Telemetry
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 3, "Target Construction & Telemetry",
              "Building the differential dataset")


# ==========================================================================
# SLIDE 10 — CVE Selection & Targets
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "CVE Selection & Target Construction")

data = [
    ["Target",       "libxml2 Tag", "CVE",            "Vulnerability Class",   "Location"],
    ["xml005_buggy", "v2.9.4",      "CVE-2017-5130",  "Integer overflow",      "xmlMemStrdupLoc()"],
    ["xml005_fixed", "v2.9.5",      "(patched)",       "\u2014",                "\u2014"],
    ["xml017_buggy", "v2.9.3",      "CVE-2016-1762",  "Heap buffer overread",  "xmlNextChar()"],
    ["xml017_fixed", "v2.9.4",      "(patched)",       "\u2014",                "\u2014"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.5),
          5, 5, data, font_size=14)

tf = add_text(slide, Inches(0.7), Inches(4.3), Inches(11), Inches(3),
              "Why these two CVEs?", font_size=18, color=ACCENT, bold=True)
add_bullet(tf, "Different vulnerability classes \u2192 tests whether features capture general structural properties", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Different libxml2 subsystems (memory management vs parser internals)", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Both statically verified: buggy code contains vulnerability, fixed code contains patch", font_size=15, color=LIGHT_GREY)
add_para(tf, "", font_size=8)
add_para(tf, "Build: AFL++ instrumentation + ASAN, unmodified FuzzBench harness (byte-identical provenance)",
         font_size=15, color=GREY)


# ==========================================================================
# SLIDE 11 — Telemetry Campaign
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Telemetry Collection Campaign")

data = [
    ["Parameter",       "Value"],
    ["Targets",         "4 (2 CVEs \u00d7 buggy/fixed)"],
    ["Seeds per target","3 (from 38-file FuzzBench corpus)"],
    ["Total runs",      "24 (12 telemetry + 12 baseline)"],
    ["Duration",        "~9.7 hours per run (until saturation)"],
    ["Execs per run",   "33\u201345 million"],
    ["Log interval",    "Every 1,000 execs (17-column CSV)"],
    ["Snapshot interval","Every 10,000 execs (65KB bitmap)"],
    ["Total data",      "~3.4 GB (40,000 snapshots + CSVs)"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(4.5),
          9, 2, data, font_size=14)

tf = add_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5),
              "Telemetry mutator design:", font_size=18, color=ACCENT, bold=True)
add_bullet(tf, "Uniform random mutation selection (no RL)", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Isolates data collection from learning effects", font_size=15, color=LIGHT_GREY)
add_para(tf, "", font_size=8)
add_para(tf, "17 metrics per step:", font_size=16, color=WHITE, bold=True)
add_bullet(tf, "Coverage: total_edges, new_edges, discovery_rate", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Edge heat: hot/warm/cool/cold classification", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Distribution: entropy, hit_mean, hit_std, hit_max", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Execution: avg_exec_time, corpus_size, crashes", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Cumulative bitmap via max-merge", font_size=14, color=LIGHT_GREY)
add_para(tf, "", font_size=8)
add_para(tf, "Saturation results:", font_size=16, color=WHITE, bold=True)
add_bullet(tf, "xml005: 5,371 \u00b1 26 (buggy) vs 5,165 \u00b1 9 (fixed)", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "xml017: 5,784 \u00b1 32 (buggy) vs 5,488 \u00b1 41 (fixed)", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Buggy consistently discovers more edges (+4\u20135%)", font_size=14, color=GREEN)


# ==========================================================================
# SLIDE 12 — Section: Differential Analysis
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 4, "Differential Analysis & Feature Derivation",
              "From raw telemetry to 13-dimensional state vector")


# ==========================================================================
# SLIDE 13 — Statistical Methods
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Statistical Methods")

tf = add_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.5),
              "Mann-Whitney U Test", font_size=20, color=ACCENT, bold=True)
add_bullet(tf, "Non-parametric comparison of buggy vs fixed", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "With n=3 per group: min p-value = 0.05", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Bonferroni threshold: \u03b1 = 0.05/65 \u2248 0.00077", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Nothing passes \u2192 expected, not a flaw", font_size=15, color=ORANGE)
add_para(tf, "", font_size=10)
add_para(tf, "Vargha-Delaney A12 Effect Size", font_size=20, color=ACCENT, bold=True)
add_bullet(tf, "A12 = P(X_buggy > X_fixed) + 0.5 \u00b7 P(equal)", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "|A12 - 0.5| \u2265 0.21: large effect", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "|A12 - 0.5| \u2265 0.14: medium effect", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Does not depend on sample size", font_size=15, color=GREEN)

tf2 = add_text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.5),
               "Why A12 over p-values?", font_size=20, color=ACCENT, bold=True)
add_bullet(tf2, "p-value answers: \"Are we confident it differs?\"", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "A12 answers: \"How often does it differ?\"", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "For feature selection, magnitude > certainty", font_size=15, color=GREEN)
add_para(tf2, "", font_size=10)
add_para(tf2, "Divergence Detection", font_size=20, color=ACCENT, bold=True)
add_bullet(tf2, "Interpolate curves to 500 common exec points", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Pooled std: \u221a(((n_b-1)\u00b7\u03c3_b\u00b2 + (n_f-1)\u00b7\u03c3_f\u00b2) / (n_b+n_f-2))", font_size=14, color=LIGHT_GREY)
add_bullet(tf2, "Divergence: |mean_b - mean_f| > pooled_std for \u22655 consecutive points", font_size=14, color=LIGHT_GREY)


# ==========================================================================
# SLIDE 14 — Key Findings
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Differential Analysis: Key Findings")

# Divergence table
data = [
    ["CVE Pair", "Divergence Point", "Buggy Final", "Fixed Final", "\u0394 Edges"],
    ["xml005 (integer overflow)",  "3,486 execs (early)",   "5,371 \u00b1 26", "5,165 \u00b1 9",  "+206"],
    ["xml017 (heap overread)",     "238,453 execs (late)",  "5,784 \u00b1 32", "5,488 \u00b1 41", "+296"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5),
          3, 5, data, font_size=14)

tf = add_text(slide, Inches(0.7), Inches(3.4), Inches(5.5), Inches(3.5),
              "Coverage divergence:", font_size=17, color=ACCENT, bold=True)
add_bullet(tf, "Buggy versions discover more edges (+4\u20135%)", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "Bug-adjacent code creates additional reachable paths", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "xml005 diverges immediately (overflow on hot path)", font_size=14, color=LIGHT_GREY)
add_bullet(tf, "xml017 diverges late (needs specific UTF-8 sequences)", font_size=14, color=LIGHT_GREY)

tf2 = add_text(slide, Inches(6.8), Inches(3.4), Inches(5.5), Inches(3.5),
               "Mutation effectiveness differs by bug class:", font_size=17, color=ACCENT, bold=True)
add_bullet(tf2, "xml005: Arithmetic mutations dominate", font_size=14, color=LIGHT_GREY)
add_bullet(tf2, "  ARITH_SUB4LE (2.2\u00d7), HAVOC_ARITH16BE (2.1\u00d7)", font_size=13, color=GREY, level=1)
add_bullet(tf2, "xml017: Dictionary/structural mutations dominate", font_size=14, color=LIGHT_GREY)
add_bullet(tf2, "  HAVOC_INT32 (1.8\u00d7), FLIP_2BITS (1.6\u00d7)", font_size=13, color=GREY, level=1)
add_para(tf2, "", font_size=6)
add_para(tf2, "Crash differential:", font_size=17, color=ACCENT, bold=True)
add_bullet(tf2, "xml005: 278 vs 0 (clean signal)", font_size=14, color=GREEN)
add_bullet(tf2, "xml017: 97 vs 87 (noisy \u2014 ASAN catches non-CVE issues)", font_size=14, color=ORANGE)


# ==========================================================================
# SLIDE 15 — Feature Ranking
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "13-Feature State Vector: Ranked by Effect Size")

data = [
    ["Rank", "Feature",           "A12 Dev", "Category",   "Key Insight"],
    ["1",    "total_edges",       "0.389",   "Coverage",   "Strongest discriminator"],
    ["2",    "cold_edges",        "0.267",   "Frontier",   "Unexplored code remaining"],
    ["3",    "corpus_size",       "0.244",   "Productivity","Interesting inputs found"],
    ["4",    "hot_edges",         "0.244",   "Heat",       "Heavily-exercised concentration"],
    ["5",    "cool_edges",        "0.233",   "Heat",       "Lightly-touched discovery edge"],
    ["6",    "avg_exec_time",     "0.222",   "Timing",     "Bug-adjacent timing anomalies"],
    ["7",    "edge_hit_mean",     "0.211",   "Depth",      "Average execution depth"],
    ["8",    "warm_edges",        "0.200",   "Heat",       "Transition zone branches"],
    ["9",    "edge_hit_std",      "0.200",   "Distribution","Execution profile variance"],
    ["10",   "edge_entropy",      "0.189",   "Distribution","Coverage shape (Shannon)"],
    ["11",   "crashes",           "0.133",   "Reward",     "Direct but noisy bug signal"],
    ["12",   "new_edges",         "0.0*",    "Reward",     "Essential RL feedback signal"],
    ["13",   "coverage_velocity", "0.0*",    "Temporal",   "Exploration vs exploitation"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5),
          14, 5, data, font_size=12)

add_text(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.4),
         "* Zero discriminative power but included as essential RL learning signals (not meant to distinguish buggy/fixed)",
         font_size=11, color=GREY)


# ==========================================================================
# SLIDE 16 — Generalization Argument
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Why These Features Generalize")

tf = add_text(slide, Inches(0.7), Inches(1.6), Inches(11), Inches(5.5),
              "Features capture structural properties of code exploration, not target-specific patterns:",
              font_size=18, color=WHITE, bold=True)
add_para(tf, "", font_size=8)

items = [
    ("Edge heat distribution (hot/warm/cool/cold ratios)",
     "Any vulnerability creates reachable code paths that shift the heat distribution \u2014 "
     "more error handlers, fallthrough cases, reachable code"),
    ("Shannon entropy",
     "Compact summary of how evenly execution effort is distributed. "
     "Bug-adjacent code creates \"hot spots\" that reduce entropy"),
    ("Execution time (EMA)",
     "Vulnerability-adjacent code involves additional processing "
     "(error handling, memory ops) creating measurable timing anomalies"),
    ("Coverage velocity (ring buffer)",
     "Rate of discovery changes as the fuzzer explores different regions "
     "of the code graph \u2014 temporal exploration signal"),
]
for title, desc in items:
    add_para(tf, title, font_size=17, color=GREEN, bold=True, space_before=Pt(14))
    add_bullet(tf, desc, font_size=14, color=LIGHT_GREY)

add_para(tf, "", font_size=8)
add_para(tf, "Tested via transfer evaluation: train on xml005 (integer overflow), evaluate on xml017 (heap overread)",
         font_size=16, color=ORANGE, bold=True)


# ==========================================================================
# SLIDE 17 — Section: Implementation
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 5, "M3_0 Implementation",
              "From feature spec to running code")


# ==========================================================================
# SLIDE 18 — SHM Layout & Normalization
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "SHM Layout (128 bytes) & Normalization Split")

data = [
    ["Offset", "Field",             "Type",    "Normalized In", "Normalization"],
    ["0",      "state_seq",         "uint32",  "\u2014",        "Sequence counter"],
    ["4",      "total_edges",       "uint32",  "Python",        "/ 65536"],
    ["8",      "cold_edges",        "uint32",  "Python",        "/ 65536"],
    ["12",     "hot_edges",         "uint32",  "Python",        "/ max(total, 1)"],
    ["16",     "warm_edges",        "uint32",  "Python",        "/ max(total, 1)"],
    ["20",     "cool_edges",        "uint32",  "Python",        "/ max(total, 1)"],
    ["24",     "edge_entropy",      "float32", "C",             "/ 3.0 (max log\u2082 8)"],
    ["28",     "edge_hit_mean",     "float32", "C",             "/ 255.0"],
    ["32",     "edge_hit_std",      "float32", "C",             "/ 255.0"],
    ["36",     "corpus_size",       "uint32",  "Python",        "log1p / log1p(10000)"],
    ["40",     "crashes",           "uint32",  "Python",        "log1p / log1p(1000)"],
    ["44",     "new_edges",         "uint32",  "Python",        "min(n,100) / 100"],
    ["48",     "avg_exec_time",     "float32", "C",             "log1p / log1p(100000)"],
    ["52",     "coverage_velocity", "float32", "C",             "min(v/0.1, 1.0)"],
    ["64",     "action_seq",        "uint32",  "\u2014",        "Sequence counter"],
    ["68",     "action",            "int32",   "\u2014",        "0\u201346 (47 mutations)"],
]
add_table(slide, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.8),
          17, 5, data, font_size=11)


# ==========================================================================
# SLIDE 19 — Section: Results
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 6, "Training & Evaluation Results", "")


# ==========================================================================
# SLIDE 20 — Eval Protocol
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Evaluation Protocol")

data = [
    ["Variant",      "Model",  "Algorithm",       "Training Target",  "Training Steps"],
    ["M3_0 DQN",     "M3_0",   "Double DQN",      "xml005_buggy",     "500,000"],
    ["M3_0 Bandit",  "M3_0",   "Thompson Sampling","xml005_buggy",     "500,000"],
    ["M1_0",         "M1_0",   "Double DQN",      "xml005_buggy",     "500,000"],
    ["Baseline",     "\u2014", "AFL++ default",    "\u2014",           "\u2014"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.5),
          5, 5, data, font_size=14)

tf = add_text(slide, Inches(0.7), Inches(4.4), Inches(11), Inches(2.5),
              "Evaluation:", font_size=18, color=ACCENT, bold=True)
add_bullet(tf, "5 runs per variant per target (for statistical significance)", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "500,000 eval steps with frozen policy (\u03b5 = 0.01)", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "xml005_buggy: in-distribution (same target as training)", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "xml017_buggy: transfer (different CVE, tests generalization)", font_size=16, color=LIGHT_GREY)
add_bullet(tf, "Plateau detection disabled \u2192 all variants train exactly 500K steps (fair comparison)", font_size=16, color=LIGHT_GREY)


# ==========================================================================
# SLIDE 21 — Results: xml005
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Results: xml005_buggy (In-Distribution)")

data = [
    ["Variant",       "Run 1", "Run 2", "Run 3", "Run 4", "Run 5", "Mean \u00b1 Std"],
    ["M3_0 DQN",      "3,945", "3,940", "4,010", "3,927", "3,962", "3,957 \u00b1 31"],
    ["M3_0 Bandit",   "3,629", "3,677", "3,644", "3,674", "3,633", "3,651 \u00b1 22"],
    ["M1_0 (DQN)",    "3,608", "3,631", "3,570", "3,528", "3,591", "3,586 \u00b1 38"],
    ["Baseline AFL++","4,320", "4,030", "4,324", "4,273", "4,302", "4,250 \u00b1 126"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(2.5),
          5, 7, data, font_size=14)

# Chart
chart_data = CategoryChartData()
chart_data.categories = ['M3_0\nDQN', 'M3_0\nBandit', 'M1_0', 'Baseline\nAFL++']
chart_data.add_series('Mean Coverage (edges)', (3957, 3651, 3586, 4250))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.5), Inches(4.2), Inches(10), Inches(3.0),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80

series = plot.series[0]
# Color individual points
for idx, color in enumerate([CHART_DQN, CHART_BAN, CHART_M10, CHART_BL]):
    pt = series.points[idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
chart.value_axis.minimum_scale = 3200
chart.value_axis.has_title = False
chart.category_axis.has_title = False


# ==========================================================================
# SLIDE 22 — Results: xml017
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Results: xml017_buggy (Transfer)")

data = [
    ["Variant",       "Run 1", "Run 2", "Run 3", "Run 4", "Run 5", "Mean \u00b1 Std"],
    ["M3_0 DQN",      "3,910", "3,876", "3,905", "3,913", "3,893", "3,899 \u00b1 15"],
    ["M3_0 Bandit",   "3,603", "3,588", "3,574", "3,601", "3,625", "3,598 \u00b1 19"],
    ["M1_0 (DQN)",    "3,588", "3,577", "3,540", "3,557", "3,572", "3,567 \u00b1 18"],
    ["Baseline AFL++","3,973", "4,303", "4,298", "4,314", "3,982", "4,174 \u00b1 173"],
]
add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(2.5),
          5, 7, data, font_size=14)

# Chart
chart_data = CategoryChartData()
chart_data.categories = ['M3_0\nDQN', 'M3_0\nBandit', 'M1_0', 'Baseline\nAFL++']
chart_data.add_series('Mean Coverage (edges)', (3899, 3598, 3567, 4174))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.5), Inches(4.2), Inches(10), Inches(3.0),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80

series = plot.series[0]
for idx, color in enumerate([CHART_DQN, CHART_BAN, CHART_M10, CHART_BL]):
    pt = series.points[idx]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
chart.value_axis.minimum_scale = 3200
chart.value_axis.has_title = False
chart.category_axis.has_title = False


# ==========================================================================
# SLIDE 23 — Pairwise Comparison
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Pairwise Comparisons & Variance")

data = [
    ["Comparison",            "xml005 \u0394", "xml005 %",  "xml017 \u0394", "xml017 %"],
    ["M3_0 DQN vs M1_0",     "+371",          "+10.3%",    "+332",          "+9.3%"],
    ["M3_0 Bandit vs M1_0",  "+65",           "+1.8%",     "+31",           "+0.9%"],
    ["M3_0 DQN vs Baseline", "\u2212293",     "\u22126.9%","\u2212275",     "\u22126.6%"],
]
add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.2),
          4, 5, data, font_size=15)

# Variance table
data2 = [
    ["Variant",       "xml005 CV", "xml017 CV", "Interpretation"],
    ["M3_0 DQN",      "0.79%",     "0.38%",     "Most stable"],
    ["M3_0 Bandit",   "0.60%",     "0.53%",     "Stable"],
    ["M1_0",          "1.07%",     "0.51%",     "Moderate"],
    ["Baseline AFL++","2.96%",     "4.15%",     "Most variable"],
]
add_table(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5),
          5, 4, data2, font_size=14)

add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.6),
         "CV = coefficient of variation (std/mean). RL models produce far more consistent results than the stochastic baseline.",
         font_size=13, color=GREY)


# ==========================================================================
# SLIDE 24 — Section: Analysis
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 7, "Analysis & Discussion", "")


# ==========================================================================
# SLIDE 25 — Key Takeaways
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Key Takeaways")

findings = [
    ("1. Differential features outperform hand-designed features",
     GREEN,
     ["+10.3% over M1_0 on in-distribution, +9.3% on transfer",
      "Consistent across all runs (no confidence interval overlap)",
      "Validates: empirically-derived > intuition-designed"]),
    ("2. Transfer generalization holds",
     GREEN,
     ["9.3% improvement on unseen xml017 (only slightly less than 10.3% on xml005)",
      "Structural features (heat, entropy, timing) are not target-specific"]),
    ("3. Policy collapse remains the bottleneck",
     RED,
     ["All RL variants underperform vanilla AFL++ by 7\u201314%",
      "DQN converges to single dominant action during eval (action 10)",
      "Loses mutation diversity that AFL++'s heuristic scheduling maintains"]),
    ("4. DQN > Bandit in sparse-reward regime",
     ORANGE,
     ["DQN: 3,957 / 3,899  vs  Bandit: 3,651 / 3,598",
      "Bandit lacks temporal credit assignment (treats steps independently)",
      "Thompson sampling exploration < \u03b5-greedy diversity in this setting"]),
]

y = 1.5
for title, color, bullets in findings:
    add_text(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    y += 0.35
    for b in bullets:
        add_text(slide, Inches(1.3), Inches(y), Inches(11), Inches(0.35),
                 "\u2022  " + b, font_size=14, color=LIGHT_GREY)
        y += 0.3
    y += 0.15


# ==========================================================================
# SLIDE 26 — Policy Collapse Deep Dive
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Why Does Baseline Still Win? Policy Collapse")

tf = add_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.5),
              "The problem:", font_size=20, color=RED, bold=True)
add_bullet(tf, "DQN learns to favor action 10 (ARITH_SUB2LE)", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "During eval (\u03b5=0.01), picks this action ~90% of the time", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Loses the mutation diversity that finds new coverage", font_size=15, color=LIGHT_GREY)
add_para(tf, "", font_size=10)
add_para(tf, "Root causes:", font_size=20, color=ORANGE, bold=True)
add_bullet(tf, "Sparse rewards: new edges become rare at saturation", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Credit assignment: \u03b3=0.99 attributes reward to recent actions, but fuzzing causality spans many executions", font_size=15, color=LIGHT_GREY)
add_bullet(tf, "Greedy eval: \u03b5=0.01 near-deterministic action selection", font_size=15, color=LIGHT_GREY)

tf2 = add_text(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.5),
               "What AFL++ does differently:", font_size=20, color=ACCENT, bold=True)
add_bullet(tf2, "Power schedules (explore/exploit scheduling)", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Queue culling (favors small, fast inputs)", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Deterministic + havoc stages", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Splice mutations across corpus", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Maintains diversity throughout campaign", font_size=15, color=GREEN)
add_para(tf2, "", font_size=10)
add_para(tf2, "Key insight:", font_size=20, color=WHITE, bold=True)
add_bullet(tf2, "The RL agent replaces ALL of this with a single argmax over 47 actions", font_size=15, color=LIGHT_GREY)
add_bullet(tf2, "Better approach: modulate AFL++'s scheduler, not replace it", font_size=15, color=GREEN)


# ==========================================================================
# SLIDE 27 — Next Steps
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
section_slide(slide, 8, "Next Steps", "")


# ==========================================================================
# SLIDE 28 — Next Steps Detail
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
title_bar(slide, "Proposed Next Steps")

items = [
    ("Address policy collapse", ACCENT,
     ["Entropy regularization: add entropy bonus to reward",
      "Action diversity constraints: minimum usage rate per action",
      "Higher eval \u03b5 (e.g., 0.1) to maintain exploration"]),
    ("Hybrid scheduling (M4 direction)", GREEN,
     ["Use RL to modulate AFL++'s existing scheduler, not replace it",
      "Agent adjusts power schedules, mutation stage weights, queue priority",
      "Preserves AFL++'s proven heuristics while adding learned adaptation"]),
    ("Better credit assignment", ORANGE,
     ["Reward shaping based on coverage velocity / entropy changes",
      "Intrinsic motivation (curiosity-driven exploration)",
      "Longer-horizon algorithms (PPO, reward redistribution)"]),
    ("Broader validation", ACCENT2,
     ["Train on multiple targets simultaneously",
      "Test on completely different software (not just libxml2)",
      "Larger seed budget (n > 3) for stronger statistical claims"]),
]

y = 1.5
for title, color, bullets in items:
    add_text(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.4),
             title, font_size=19, color=color, bold=True)
    y += 0.4
    for b in bullets:
        add_text(slide, Inches(1.3), Inches(y), Inches(11), Inches(0.35),
                 "\u2022  " + b, font_size=15, color=LIGHT_GREY)
        y += 0.32
    y += 0.15


# ==========================================================================
# SLIDE 29 — Summary
# ==========================================================================
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Summary", font_size=40, color=ACCENT, bold=True)

shape = slide.shapes.add_shape(1, Inches(1), Inches(2.0), Inches(4), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

summary = [
    ("Built", "a differential fuzzing pipeline: 4 libxml2 targets, 24 campaigns, 3.4 GB telemetry"),
    ("Derived", "13-feature state vector from buggy vs fixed comparison using A12 effect size ranking"),
    ("Implemented", "M3_0 mutator (C) + model (Python) with 128-byte SHM protocol"),
    ("Trained", "DQN and contextual bandit variants, 500K steps each"),
    ("Result", "M3_0 DQN achieves +10.3% over M1_0 (in-dist) and +9.3% (transfer)"),
    ("Gap", "Baseline AFL++ still leads by ~7% due to policy collapse"),
    ("Next", "Address collapse via hybrid scheduling / entropy regularization"),
]

y = 2.4
for label, desc in summary:
    add_text(slide, Inches(1), Inches(y), Inches(1.8), Inches(0.4),
             label, font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(slide, Inches(3.0), Inches(y), Inches(9), Inches(0.4),
             desc, font_size=18, color=WHITE)
    y += 0.55

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Questions?", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ==========================================================================
# Save
# ==========================================================================
out_path = os.path.join(os.path.dirname(__file__), "experiment_3_review.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
